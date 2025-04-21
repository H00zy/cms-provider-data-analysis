import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import os

# Paths
data_path = r"/Data/cleaned/master_df.csv"
output_img = r"C:\Users\Admin\PycharmProjects\CMS provider data analysis\outputs\mips_by_specialty.png"
ppt_path = r"/outputs/MIPS_By_Specialty_Summary.pptx"
os.makedirs(os.path.dirname(output_img), exist_ok=True)

# Load data
df = pd.read_csv(data_path)

# Filter
plot_df = df[['pri_spec', 'final_mips_score']].dropna()
plot_df = plot_df[plot_df['final_mips_score'].between(0, 100)]  # MIPS scores range from 0-100

# Keep top 15 specialties by count for readability
top_specs = plot_df['pri_spec'].value_counts().nlargest(15).index.tolist()
plot_df = plot_df[plot_df['pri_spec'].isin(top_specs)]

# Plot
plt.figure(figsize=(14, 8))
sns.boxplot(x='pri_spec', y='final_mips_score', data=plot_df)
plt.xticks(rotation=45, ha='right')
plt.title("MIPS Score Distribution by Specialty (Top 15)", fontsize=14)
plt.xlabel("Clinician Specialty")
plt.ylabel("Final MIPS Score")
plt.tight_layout()
plt.savefig(output_img)
plt.show()

print(f"✅ Saved MIPS score graph to: {output_img}")

# Create PowerPoint
prs = Presentation()
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.title
title_box.text = "MIPS Score Distribution by Specialty (Top 15)"

# Add image
left = Inches(0.5)
top = Inches(1.2)
height = Inches(5)
slide.shapes.add_picture(output_img, left, top, height=height)

# Save presentation
prs.save(ppt_path)
print(f"✅ PowerPoint slide saved to: {ppt_path}")
